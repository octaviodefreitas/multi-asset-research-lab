"""Presentation exports: PowerPoint tearsheet and Excel workbook.

PowerPoint charts are written as native chart objects (not images), so an
analyst can recolor and restyle them to a house template after pasting.
Slides are 16:9 with sparse yearly axis labels and bottom legends.
"""
from __future__ import annotations

from datetime import date
from io import BytesIO

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

DISCLAIMER = ("Backtested / simulated research on historical data. Not live trading "
              "performance; not investment advice. Backtested results are subject to "
              "overfitting and will differ from live results.")


def _sparse_year_labels(index: pd.DatetimeIndex) -> list[str]:
    """One label per year, blank otherwise — keeps the category axis readable."""
    labels, seen = [], set()
    for ts in index:
        if ts.year not in seen:
            labels.append(str(ts.year))
            seen.add(ts.year)
        else:
            labels.append("")
    return labels


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    ph = slide.placeholders[1]
    ph.text = f"{subtitle}\nGenerated {date.today():%d %b %Y}\n\n{DISCLAIMER}"
    for para in ph.text_frame.paragraphs:
        para.font.size = Pt(14)


def _chart_slide(prs: Presentation, title: str, series: dict[str, pd.Series],
                 number_format: str = "0.00", kind: str = "line") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)

    data = CategoryChartData()
    first = next(iter(series.values()))
    if kind == "bar":
        data.categories = [str(i) for i in first.index]
        chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    else:
        data.categories = _sparse_year_labels(first.index)
        chart_type = XL_CHART_TYPE.LINE
    for name, s in series.items():
        data.add_series(name, tuple(None if pd.isna(v) else round(float(v), 5)
                                    for v in s.to_numpy()))

    chart = slide.shapes.add_chart(
        chart_type, Inches(0.45), Inches(1.35), Inches(12.4), Inches(5.7), data
    ).chart
    chart.has_title = False
    chart.font.size = Pt(11)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.tick_labels.number_format = number_format
    chart.value_axis.tick_labels.number_format_is_linked = False
    if kind == "line":
        for ser in chart.series:
            ser.smooth = True
            ser.format.line.width = Pt(1.75)


def _table_slide(prs: Presentation, title: str, df: pd.DataFrame) -> None:
    """df must already contain display-formatted strings."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
    rows, cols = df.shape[0] + 1, df.shape[1] + 1
    table = slide.shapes.add_table(rows, cols, Inches(0.45), Inches(1.5),
                                   Inches(12.4), Inches(0.42 * rows)).table
    table.columns[0].width = Inches(3.0)
    table.cell(0, 0).text = ""
    for j, col in enumerate(df.columns):
        table.cell(0, j + 1).text = str(col)
    for i, (idx, row) in enumerate(df.iterrows()):
        table.cell(i + 1, 0).text = str(idx)
        for j, value in enumerate(row):
            table.cell(i + 1, j + 1).text = str(value)
    for r, cell_row in enumerate(table.rows):
        for cell in cell_row.cells:
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12 if r == 0 else 11)


def build_tearsheet(subtitle: str,
                    charts: list[tuple[str, dict[str, pd.Series], str, str]],
                    tables: list[tuple[str, pd.DataFrame]]) -> bytes:
    """Assemble a 16:9 deck: title slide, one slide per chart (native,
    editable; kind "line" or "bar"), one slide per formatted table."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _title_slide(prs, "Multi-Asset Systematic Strategy — Tearsheet", subtitle)
    for title, series, number_format, kind in charts:
        _chart_slide(prs, title, series, number_format, kind)
    for title, df in tables:
        _table_slide(prs, title, df)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_workbook(sheets: dict[str, pd.DataFrame]) -> bytes:
    """Multi-sheet Excel workbook; sheet order follows the dict order."""
    buf = BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for name, df in sheets.items():
            df.to_excel(writer, sheet_name=name[:31])
    return buf.getvalue()
