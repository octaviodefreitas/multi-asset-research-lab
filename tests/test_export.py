from io import BytesIO

import numpy as np
import pandas as pd
from pptx import Presentation

from app import export


def _series(n=100):
    idx = pd.bdate_range("2020-01-01", periods=n)
    return pd.Series(np.linspace(1.0, 1.5, n), index=idx)


def test_tearsheet_is_valid_pptx_with_expected_slides():
    yearly = pd.Series([0.05, -0.12, 0.20], index=[2020, 2021, 2022])
    charts = [
        ("Equity", {"Strategy": _series(), "Benchmark": _series() * 0.9}, "0.00", "line"),
        ("Yearly", {"Strategy": yearly}, "0%", "bar"),
    ]
    tables = [("Metrics", pd.DataFrame({"Sharpe": ["1.10"], "CAGR": ["8.0%"]},
                                       index=["EW Portfolio"]))]
    blob = export.build_tearsheet("test subtitle", charts, tables)

    assert blob[:2] == b"PK"  # pptx is a zip container
    prs = Presentation(BytesIO(blob))
    assert len(prs.slides) == 4  # title + 2 charts + 1 table
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)  # widescreen


def test_tearsheet_handles_nans_in_series():
    s = _series()
    s.iloc[10:20] = np.nan
    blob = export.build_tearsheet("sub", [("Equity", {"Strategy": s}, "0.00", "line")], [])
    assert blob[:2] == b"PK"


def test_disclaimer_on_title_slide():
    blob = export.build_tearsheet("sub", [], [])
    prs = Presentation(BytesIO(blob))
    text = " ".join(shape.text_frame.text for shape in prs.slides[0].shapes
                    if shape.has_text_frame)
    assert "not investment advice" in text.lower()


def test_sparse_year_labels():
    idx = pd.bdate_range("2020-11-01", "2022-03-01", freq="ME")
    labels = export._sparse_year_labels(idx)
    assert labels.count("2020") == 1 and labels.count("2021") == 1
    assert labels.count("") == len(labels) - 3


def test_workbook_round_trip():
    sheets = {
        "Key Metrics": pd.DataFrame({"Sharpe": [1.1]}, index=["EW"]),
        "Equity Curves": pd.DataFrame({"Strategy": _series()}),
    }
    blob = export.build_workbook(sheets)
    assert blob[:2] == b"PK"
    back = pd.read_excel(BytesIO(blob), sheet_name=None)
    assert set(back) == {"Key Metrics", "Equity Curves"}
    assert back["Key Metrics"]["Sharpe"].iloc[0] == 1.1
